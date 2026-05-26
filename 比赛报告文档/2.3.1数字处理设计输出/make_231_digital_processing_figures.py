from __future__ import annotations

import math
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
from PIL import Image
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Rectangle
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor


BASE_DIR = Path(__file__).resolve().parent
FIG_DIR = BASE_DIR / "figures"
DOCX_PATH = BASE_DIR / "2.3.1 数字处理设计.docx"
PPTX_PATH = BASE_DIR / "2.3.1 数字处理设计插图包.pptx"

mpl.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": [
        "Microsoft YaHei",
        "SimHei",
        "Arial Unicode MS",
        "DejaVu Sans",
        "sans-serif",
    ],
    "axes.unicode_minus": False,
    "svg.fonttype": "none",
    "pdf.fonttype": 42,
    "font.size": 9,
    "axes.linewidth": 0.8,
    "legend.frameon": False,
})

PALETTE = {
    "ink": "#263238",
    "muted": "#667085",
    "blue": "#5B8DEF",
    "blue2": "#DCE9FF",
    "teal": "#52B6B0",
    "teal2": "#DDF3F1",
    "green": "#5BAE6B",
    "green2": "#E3F2E7",
    "amber": "#D99A3D",
    "amber2": "#F8EBD7",
    "rose": "#D8737F",
    "rose2": "#F6DFE3",
    "violet": "#8B7BD3",
    "violet2": "#E9E5FA",
    "gray": "#EEF1F4",
    "line": "#B7C0C8",
}

RBW_ROWS = [
    ("1 kHz", 13000, 5, 256, 256, 128),
    ("10 kHz", 1300, 5, 256, 256, 128),
    ("30 kHz", 433, 5, 256, 256, 128),
    ("100 kHz", 130, 5, 128, 384, 64),
    ("300 kHz", 43, 4, 128, 384, 64),
    ("1 MHz", 13, 4, 64, 384, 32),
]

FIGURES = [
    {
        "id": "2.3.1-1",
        "file": "fig_231_1_overall_flow",
        "title": "Zynq 数字中频处理总体流程",
        "caption": "图 2.3.1-1 Zynq 数字中频处理总体流程。40 MHz 中频经 ADC 数字化后，由 DMA 搬运到 PS 侧缓存，随后完成 DDC、RBW 滤波、功率估计和幅度校正，并通过串口协议输出频谱点。",
        "slide_note": "从固定中频到频谱点的主处理链路。",
    },
    {
        "id": "2.3.1-2",
        "file": "fig_231_2_software_arch",
        "title": "Zynq 内部软件与数据流架构",
        "caption": "图 2.3.1-2 Zynq 内部软件与数据流架构。扫频状态机作为调度核心，协调本振控制、DMA 采集、信号处理和协议输出，使每个频点的采样、测量、发送过程解耦。",
        "slide_note": "突出 sweep_engine、signal_processing、dma_capture 和 device_protocol 的协同关系。",
    },
    {
        "id": "2.3.1-3",
        "file": "fig_231_3_sweep_state",
        "title": "逐点扫频状态机",
        "caption": "图 2.3.1-3 逐点扫频状态机。每个频点依次完成 LO1 设置、锁定等待、DMA 采样、样本累积、功率测量、频谱点发送和下一频点切换。",
        "slide_note": "说明本系统采用频谱仪式逐点扫频，而不是一次 FFT 上传。",
    },
    {
        "id": "2.3.1-4",
        "file": "fig_231_4_rbw_chain",
        "title": "多档 RBW 滤波处理链",
        "caption": "图 2.3.1-4 多档 RBW 滤波处理链。不同 RBW 档位配置不同 CIC 抽取比、CIC 级数和补偿 FIR 长度，从而在扫速和分辨率之间取得折中。",
        "slide_note": "展示 RBW 档位如何映射为数字滤波和抽取参数。",
    },
    {
        "id": "2.3.1-5",
        "file": "fig_231_5_power_calibration",
        "title": "功率计算与幅度校正路径",
        "caption": "图 2.3.1-5 功率计算与幅度校正路径。RBW 滤波后的 I/Q 数据先通过均方功率估计得到 dBFS，再结合 ADC 满量程参考和射频前端校正换算为输入端 dBm。",
        "slide_note": "把测量数学关系和前端增益/衰减补偿放在同一条链上。",
    },
    {
        "id": "2.3.1-表1",
        "file": "table_231_1_rbw_params",
        "title": "RBW 档位参数表",
        "caption": "表 2.3.1-1 RBW 档位参数。参数来自当前固件 app_config.h，用于配置 CIC 抽取、补偿 FIR 和每个扫频点的有效观测长度。",
        "slide_note": "当前固件中各 RBW 档位的关键数字处理参数。",
    },
]


def ensure_dirs() -> None:
    FIG_DIR.mkdir(parents=True, exist_ok=True)


def add_box(ax, xy, w, h, text, fc, ec=None, fontsize=10, radius=0.04, weight="normal"):
    if ec is None:
        ec = PALETTE["line"]
    patch = FancyBboxPatch(
        xy,
        w,
        h,
        boxstyle=f"round,pad=0.018,rounding_size={radius}",
        facecolor=fc,
        edgecolor=ec,
        linewidth=1.2,
    )
    ax.add_patch(patch)
    ax.text(
        xy[0] + w / 2,
        xy[1] + h / 2,
        text,
        ha="center",
        va="center",
        fontsize=fontsize,
        color=PALETTE["ink"],
        fontweight=weight,
        linespacing=1.25,
    )
    return patch


def add_arrow(ax, start, end, color=None, lw=1.8, style="-|>", rad=0.0):
    if color is None:
        color = PALETTE["muted"]
    arrow = FancyArrowPatch(
        start,
        end,
        arrowstyle=style,
        mutation_scale=12,
        linewidth=lw,
        color=color,
        connectionstyle=f"arc3,rad={rad}",
    )
    ax.add_patch(arrow)
    return arrow


def add_title(ax, text, subtitle=None):
    ax.text(0.02, 0.96, text, ha="left", va="top", fontsize=14, fontweight="bold", color=PALETTE["ink"])
    if subtitle:
        ax.text(0.02, 0.905, subtitle, ha="left", va="top", fontsize=8.5, color=PALETTE["muted"])


def save_figure(fig, name):
    out = FIG_DIR / name
    fig.savefig(out.with_suffix(".png"), dpi=300, bbox_inches="tight", facecolor="white")
    fig.savefig(out.with_suffix(".svg"), bbox_inches="tight", facecolor="white")
    fig.savefig(out.with_suffix(".pdf"), bbox_inches="tight", facecolor="white")
    fig.savefig(out.with_suffix(".tiff"), dpi=600, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def setup_canvas(width=11.5, height=6.2):
    fig, ax = plt.subplots(figsize=(width, height))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    return fig, ax


def draw_overall_flow():
    fig, ax = setup_canvas()
    add_title(
        ax,
        "数字中频到频谱点的主处理链路",
        "Core conclusion: Zynq 将 40 MHz 中频转换为带 RBW 选择和幅度校正的定量频谱点。",
    )
    xs = [0.035, 0.175, 0.315, 0.455, 0.595, 0.735, 0.875]
    box_w = 0.098
    labels = [
        "LTC2208\nADC采样\n130 MSPS",
        "AXI DMA\n4096点帧\nPS缓存",
        "DDC\n40 MHz IF\nI/Q基带",
        "CIC抽取\n按RBW降采样",
        "补偿FIR\nRBW低通\n去瞬态",
        "功率估计\nmean(I²+Q²)\ndBFS/dBm",
        "幅度校正\nUART流式\n频谱点",
    ]
    colors = [PALETTE["blue2"], PALETTE["blue2"], PALETTE["teal2"], PALETTE["teal2"], PALETTE["green2"], PALETTE["amber2"], PALETTE["rose2"]]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_box(ax, (x, 0.47), box_w, 0.19, label, color, fontsize=8.2)
        if i < len(xs) - 1:
            add_arrow(ax, (x + box_w + 0.006, 0.565), (xs[i + 1] - 0.009, 0.565))
    ax.text(0.035, 0.33, "测量方式：逐频点扫频功率测量，而非一次性上传 FFT bin", ha="left", fontsize=10.5, color=PALETTE["ink"], fontweight="bold")
    add_box(ax, (0.035, 0.18), 0.28, 0.09, "本振扫频决定当前 RF 频点\nLO1 = RF + 2.180 GHz，LO2 = 2.220 GHz", PALETTE["gray"], fontsize=8.5)
    add_box(ax, (0.36, 0.18), 0.28, 0.09, "RBW 决定内部步进与采样积累\nstep ≈ RBW / 2，最多 4096 点", PALETTE["gray"], fontsize=8.5)
    add_box(ax, (0.685, 0.18), 0.28, 0.09, "输出为频率-幅度点\n上位机实时拼接并显示频谱曲线", PALETTE["gray"], fontsize=8.5)
    save_figure(fig, "fig_231_1_overall_flow")


def draw_software_arch():
    fig, ax = setup_canvas()
    add_title(ax, "Zynq 侧软件模块协同关系", "Bare-metal 固件通过函数回调和状态机组织采集、处理与发送。")
    add_box(ax, (0.39, 0.68), 0.22, 0.12, "sweep_engine\n扫频状态机\n调度每个频点", PALETTE["violet2"], PALETTE["violet"], fontsize=9.5, weight="bold")
    modules = [
        ((0.07, 0.68), "lo_control\n设置 LO1\n等待锁定", PALETTE["blue2"], PALETTE["blue"]),
        ((0.07, 0.42), "dma_capture\nAXI DMA S2MM\n采集中频帧", PALETTE["teal2"], PALETTE["teal"]),
        ((0.39, 0.42), "signal_processing\nDDC + CIC + FIR\n功率估计", PALETTE["green2"], PALETTE["green"]),
        ((0.71, 0.42), "amplitude_correction\n衰减/LNA/VGA/频率\n输入端幅度修正", PALETTE["amber2"], PALETTE["amber"]),
        ((0.71, 0.68), "device_protocol\n0x82频谱点\nACK/状态/性能数据", PALETTE["rose2"], PALETTE["rose"]),
    ]
    for xy, text, fc, ec in modules:
        add_box(ax, xy, 0.22, 0.12, text, fc, ec, fontsize=8.8)
    add_arrow(ax, (0.29, 0.74), (0.39, 0.74), PALETTE["blue"])
    add_arrow(ax, (0.50, 0.68), (0.18, 0.54), PALETTE["teal"], rad=0.12)
    add_arrow(ax, (0.29, 0.48), (0.39, 0.48), PALETTE["teal"])
    add_arrow(ax, (0.61, 0.48), (0.71, 0.48), PALETTE["green"])
    add_arrow(ax, (0.82, 0.54), (0.82, 0.68), PALETTE["rose"])
    add_arrow(ax, (0.71, 0.74), (0.61, 0.74), PALETTE["rose"], style="<|-", rad=0.0)
    add_box(ax, (0.19, 0.18), 0.62, 0.12, "主循环持续轮询 device_protocol 和 sweep_engine；扫频活跃时释放后台 FFT/状态采集路径，保证 DMA 资源被当前测量点独占。", PALETTE["gray"], fontsize=9.2)
    save_figure(fig, "fig_231_2_software_arch")


def draw_sweep_state():
    fig, ax = setup_canvas()
    add_title(ax, "逐点扫频状态机", "每个频点都经历“锁相-采样-积累-测量-发送”的闭环。")
    states = [
        ("Prepare\n配置RBW/DDC", 0.07, 0.67, PALETTE["gray"]),
        ("Set LO1\n切到RF频点", 0.25, 0.67, PALETTE["blue2"]),
        ("Wait Lock\n等待PLL锁定", 0.43, 0.67, PALETTE["blue2"]),
        ("Arm DMA\n启动采样", 0.61, 0.67, PALETTE["teal2"]),
        ("Accumulate\nDDC+CIC累积", 0.61, 0.39, PALETTE["teal2"]),
        ("Measure\nFIR+功率估计\n幅度校正", 0.43, 0.39, PALETTE["green2"]),
        ("Emit Point\n发送0x82", 0.25, 0.39, PALETTE["rose2"]),
        ("Next Point\n或Done", 0.07, 0.39, PALETTE["amber2"]),
    ]
    for label, x, y, color in states:
        add_box(ax, (x, y), 0.13, 0.13, label, color, fontsize=8.5)
    centers = [(x + 0.065, y + 0.065) for _, x, y, _ in states]
    for i in range(3):
        add_arrow(ax, (centers[i][0] + 0.065, centers[i][1]), (centers[i + 1][0] - 0.065, centers[i + 1][1]))
    add_arrow(ax, (centers[3][0], centers[3][1] - 0.07), (centers[4][0], centers[4][1] + 0.07))
    add_arrow(ax, (centers[4][0] - 0.065, centers[4][1]), (centers[5][0] + 0.065, centers[5][1]))
    add_arrow(ax, (centers[5][0] - 0.065, centers[5][1]), (centers[6][0] + 0.065, centers[6][1]))
    add_arrow(ax, (centers[6][0] - 0.065, centers[6][1]), (centers[7][0] + 0.065, centers[7][1]))
    add_arrow(ax, (centers[7][0], centers[7][1] + 0.07), (centers[1][0] - 0.065, centers[1][1] - 0.01), PALETTE["amber"], rad=-0.36)
    add_box(ax, (0.79, 0.53), 0.16, 0.16, "异常保护\nLO锁定超时\nDMA超时\n功率测量失败", PALETTE["rose2"], PALETTE["rose"], fontsize=8.5)
    add_arrow(ax, (0.74, 0.68), (0.79, 0.63), PALETTE["rose"], rad=-0.15)
    add_arrow(ax, (0.74, 0.45), (0.79, 0.58), PALETTE["rose"], rad=0.15)
    ax.text(0.07, 0.22, "设计收益：扫频过程中每个频点独立完成测量，既便于实时显示，也便于定位 DMA、锁相或信号处理瓶颈。", fontsize=10, color=PALETTE["ink"])
    save_figure(fig, "fig_231_3_sweep_state")


def draw_rbw_chain():
    fig, ax = setup_canvas()
    add_title(ax, "RBW 多档数字滤波链", "RBW 档位决定抽取比、补偿 FIR 长度和每点观测样本数。")
    add_box(ax, (0.05, 0.56), 0.16, 0.13, "DDC后 I/Q\n130 MSPS 等效输入", PALETTE["blue2"], fontsize=9)
    add_box(ax, (0.28, 0.56), 0.16, 0.13, "CIC抽取\nR/N按RBW配置\n降低数据率", PALETTE["teal2"], fontsize=9)
    add_box(ax, (0.51, 0.56), 0.16, 0.13, "补偿FIR\nHamming窗低通\n修正通带", PALETTE["green2"], fontsize=9)
    add_box(ax, (0.74, 0.56), 0.16, 0.13, "有效观测窗\n跳过瞬态\n积累功率", PALETTE["amber2"], fontsize=9)
    for start_x in [0.21, 0.44, 0.67]:
        add_arrow(ax, (start_x, 0.625), (start_x + 0.07, 0.625))
    rbw_labels = ["1 kHz", "10 kHz", "30 kHz", "100 kHz", "300 kHz", "1 MHz"]
    rs = [13000, 1300, 433, 130, 43, 13]
    taps = [256, 256, 256, 128, 128, 64]
    x0, y0 = 0.08, 0.26
    width = 0.84
    bar_h = 0.032
    max_log = math.log10(max(rs))
    for idx, (label, r, tap) in enumerate(zip(rbw_labels, rs, taps)):
        y = y0 + (5 - idx) * 0.045
        ax.text(x0, y + bar_h / 2, label, ha="right", va="center", fontsize=8.2, color=PALETTE["ink"])
        cic_w = 0.12 + 0.22 * math.log10(r) / max_log
        ax.add_patch(Rectangle((x0 + 0.02, y), cic_w, bar_h, facecolor=PALETTE["teal"], edgecolor="none", alpha=0.85))
        ax.add_patch(Rectangle((x0 + 0.36, y), 0.10 + tap / 256 * 0.18, bar_h, facecolor=PALETTE["green"], edgecolor="none", alpha=0.8))
        ax.text(x0 + 0.03 + cic_w, y + bar_h / 2, f"R={r}", va="center", fontsize=7.5, color=PALETTE["muted"])
        ax.text(x0 + 0.66, y + bar_h / 2, f"FIR={tap} taps", va="center", fontsize=7.5, color=PALETTE["muted"])
    ax.text(0.19, 0.205, "CIC抽取比", fontsize=8.5, color=PALETTE["teal"], fontweight="bold")
    ax.text(0.48, 0.205, "补偿FIR长度", fontsize=8.5, color=PALETTE["green"], fontweight="bold")
    save_figure(fig, "fig_231_4_rbw_chain")


def draw_power_calibration():
    fig, ax = setup_canvas()
    add_title(ax, "功率估计与输入端幅度校正", "把基带样本能量转换为可读的 RF 输入端 dBm。")
    add_box(ax, (0.06, 0.62), 0.18, 0.12, "RBW滤波后\nI/Q样本", PALETTE["green2"], fontsize=9.5)
    add_box(ax, (0.31, 0.62), 0.18, 0.12, "均方功率\nP = mean(I²+Q²)", PALETTE["green2"], fontsize=9.5)
    add_box(ax, (0.56, 0.62), 0.18, 0.12, "dBFS换算\n10log10(P/PFS)", PALETTE["amber2"], fontsize=9.5)
    add_box(ax, (0.79, 0.62), 0.16, 0.12, "ADC输入功率\n+8.02 dBm参考", PALETTE["amber2"], fontsize=9.5)
    for x in [0.24, 0.49, 0.74]:
        add_arrow(ax, (x, 0.68), (x + 0.07, 0.68))
    add_box(ax, (0.23, 0.32), 0.54, 0.13, "幅度校正项 = 衰减器 - LNA增益 - VGA增益 + 路径校准 + 频率分段校准", PALETTE["rose2"], PALETTE["rose"], fontsize=10)
    add_arrow(ax, (0.87, 0.62), (0.65, 0.45), PALETTE["amber"], rad=0.18)
    add_arrow(ax, (0.50, 0.32), (0.50, 0.22), PALETTE["rose"])
    add_box(ax, (0.35, 0.10), 0.30, 0.10, "RF输入端幅度\n用于上位机频谱显示", PALETTE["blue2"], PALETTE["blue"], fontsize=10, weight="bold")
    ax.text(0.08, 0.49, "固件保留 raw power 与 correction dB，便于后续校准和误差分析。", fontsize=9.2, color=PALETTE["muted"])
    save_figure(fig, "fig_231_5_power_calibration")


def draw_rbw_table():
    fig, ax = setup_canvas(width=11.2, height=5.8)
    add_title(ax, "当前固件 RBW 档位参数", "Source: app_config.h；用于配置 signal_processing.c 中的 sweep-path。")
    columns = ["RBW", "CIC R", "CIC N", "FIR taps", "Observe pts", "Skip pts", "Accum target"]
    cell_text = []
    for rbw, r, n, taps, obs, skip in RBW_ROWS:
        cell_text.append([rbw, str(r), str(n), str(taps), str(obs), str(skip), str(obs + skip + taps)])
    table = ax.table(
        cellText=cell_text,
        colLabels=columns,
        cellLoc="center",
        colLoc="center",
        bbox=[0.06, 0.12, 0.88, 0.68],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(9)
    for (row, col), cell in table.get_celld().items():
        cell.set_edgecolor("#D0D5DD")
        cell.set_linewidth(0.8)
        if row == 0:
            cell.set_facecolor(PALETTE["blue"])
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        else:
            cell.set_facecolor("#FFFFFF" if row % 2 else "#F8FAFC")
            if col == 0:
                cell.get_text().set_fontweight("bold")
    ax.text(0.06, 0.055, "注：Accum target = Observe points + Skip points + FIR taps，实际积累上限受 ACCUM_BUFFER_SIZE = 768 约束。", fontsize=8.5, color=PALETTE["muted"])
    save_figure(fig, "table_231_1_rbw_params")


def generate_figures():
    draw_overall_flow()
    draw_software_arch()
    draw_sweep_state()
    draw_rbw_chain()
    draw_power_calibration()
    draw_rbw_table()


def set_doc_defaults(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Cm(2.2)
    section.bottom_margin = Cm(2.0)
    section.left_margin = Cm(2.4)
    section.right_margin = Cm(2.4)
    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "宋体"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    normal.font.size = Pt(10.5)
    for name in ["Heading 1", "Heading 2", "Heading 3"]:
        style = styles[name]
        style.font.name = "黑体"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "黑体")
        style.font.color.rgb = RGBColor(38, 50, 56)
    styles["Heading 1"].font.size = Pt(16)
    styles["Heading 2"].font.size = Pt(13)
    styles["Heading 3"].font.size = Pt(11.5)


def qn(tag):
    from docx.oxml.ns import qn as _qn
    return _qn(tag)


def add_doc_paragraph(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.first_line_indent = Cm(0.74)
    p.paragraph_format.line_spacing = 1.25
    p.paragraph_format.space_after = Pt(5)
    run = p.add_run(text)
    run.font.name = "宋体"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    run.font.size = Pt(10.5)


def add_doc_caption(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(8)
    run = p.add_run(text)
    run.font.name = "宋体"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(85, 96, 110)


def add_doc_figure(doc: Document, fig_meta: dict, width_cm: float = 15.5) -> None:
    image = FIG_DIR / f"{fig_meta['file']}.png"
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(image), width=Cm(width_cm))
    add_doc_caption(doc, fig_meta["caption"])


def add_word_table(doc: Document) -> None:
    doc.add_heading("表 2.3.1-1 RBW 档位参数", level=3)
    table = doc.add_table(rows=1, cols=7)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    headers = ["RBW", "CIC R", "CIC N", "FIR taps", "Observe pts", "Skip pts", "Accum target"]
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(9)
    for row in RBW_ROWS:
        rbw, r, n, taps, obs, skip = row
        values = [rbw, str(r), str(n), str(taps), str(obs), str(skip), str(obs + skip + taps)]
        cells = table.add_row().cells
        for i, value in enumerate(values):
            cells[i].text = value
            cells[i].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            for paragraph in cells[i].paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(9)
    add_doc_caption(doc, "参数来源：当前活动固件 app_config.h。Accum target 为 Observe points、Skip points 和 FIR taps 的合计值。")


def generate_word():
    doc = Document()
    set_doc_defaults(doc)
    doc.add_heading("2.3.1 数字处理设计", level=1)

    add_doc_paragraph(doc, "本仪器的射频前端通过两级超外差结构将 1 MHz-1.5 GHz 输入信号变换到固定中频，再由中频板上的高速 ADC 完成数字化。Zynq 端数字处理设计的任务，是在有限成本和有限片上资源条件下，把固定中频采样数据转化为具有分辨率带宽选择、幅度定量能力和实时显示能力的频谱曲线。与通用 SDR 软件中常见的整段 FFT 显示不同，本系统当前固件采用更接近传统扫频频谱仪的逐点测量方式：每个频点先设置第一本振，等待锁相稳定后采集中频数据，再在 Zynq 内部完成数字下变频、RBW 滤波、功率估计与幅度校正，最后把频率-幅度点流式发送到上位机。")
    add_doc_paragraph(doc, "这种设计把硬件频率规划和数字信号处理紧密结合起来。射频输入经两级混频后形成固定低中频，扫频动作由 LO1 的频率变化完成；数字处理链路始终围绕固定中频工作，因此 NCO、滤波器和功率计算模型可以保持稳定。系统主链路如图 2.3.1-1 所示，整体可概括为“ADC 采样-DMA 搬运-DDC 基带化-CIC 抽取-RBW 补偿 FIR-功率估计-幅度校正-协议输出”。")
    add_doc_figure(doc, FIGURES[0])

    doc.add_heading("1 数字中频处理总体架构", level=2)
    add_doc_paragraph(doc, "Zynq 端固件运行在裸机环境下，主要模块包括 dma_capture、signal_processing、sweep_engine、lo_control、amplitude_correction 和 device_protocol。main.c 在初始化阶段依次完成外设、射频控制、协议引擎、数字信号处理、扫频状态机和 DMA 采集模块的初始化；运行阶段由主循环轮询协议命令和扫频状态机。上位机通过 GET_SPECTRUM 或 START_SWEEP 等命令触发扫频后，device_protocol 把命令转交给 sweep_engine，后者成为整个数字处理链路的调度核心。")
    add_doc_paragraph(doc, "软件架构上的一个重要特点是采集、处理和输出之间采用明确的接口分层。dma_capture 只负责 AXI DMA S2MM 通道的启动、完成标志和错误状态；signal_processing 只负责中频数据的 DDC、抽取、滤波和功率估计；sweep_engine 负责把 LO1 控制、DMA 采样、样本积累、功率计算、幅度校正和点数据发送串联为一个非阻塞状态机；device_protocol 负责帧格式、ACK、状态数据和频谱点数据的串口发送。这种分层使得调试时可以分别定位锁相、DMA、滤波、功率估计或协议发送问题。")
    add_doc_figure(doc, FIGURES[1])

    doc.add_heading("2 ADC采样与DMA数据搬运", level=2)
    add_doc_paragraph(doc, "中频板采用 LTC2208 对模拟中频信号进行采样，当前固件中 ADC 采样率配置为 130 MSPS。ADC 输出的数据经 PL 侧接口进入 AXI DMA，再通过 S2MM 通道搬运到 PS 侧缓存。数字处理的基本数据帧长度为 4096 点，每个采样点按 16 bit 有符号数处理，进入 signal_processing 后统一归一化到 -1 到 1 附近的浮点范围。这一帧长一方面便于复用 CMSIS-DSP 中的 FIR 和 FFT 基础设施，另一方面也给 DMA 中断、缓存管理和实时轮询提供了较清晰的处理粒度。")
    add_doc_paragraph(doc, "工程实现中，DMA 传输长度并不是简单地越大越好。当前硬件导出的 AXI DMA 简单传输长度宽度约束要求单次传输保持在 65535 字节以内，固件因此对单次传输进行了保护。更重要的是，当前扫频测量路径把每次 DMA 触发保持为一个 4096 点帧，并由扫频状态机根据 RBW 档位反复重新启动 DMA，直到抽取后的累积样本数达到目标。这样做虽然增加了 rearm 次数，但能避免多帧简单传输完成后后续块无有效样本、窄 RBW 档位被 FIR 瞬态吞掉等问题，实际可靠性更高。")
    add_doc_paragraph(doc, "在系统运行时，DMA 采样和后台状态采集共享同一 DMA 资源，因此固件在扫频活跃时以当前测量点为最高优先级。扫频状态机进入采样阶段后会复位并启动 DMA，等待 frame_ready 或错误标志；如果当前 RBW 所需样本尚未积累够，则进入 REARM_DMA 状态继续采样。这样每个频点都有明确的采样闭环，便于上位机根据状态和性能数据判断扫速瓶颈。")

    doc.add_heading("3 数字下变频与I/Q基带生成", level=2)
    add_doc_paragraph(doc, "由于硬件链路把被测信号变换到固定低中频，Zynq 端首先需要通过数字下变频把中频信号搬移到基带。固件在 signal_processing_init 中建立 NCO 正弦和余弦查找表，在扫频准备阶段按实际第二中频设置 DDC 频率。当前频率规划中 LO2 固定为 2.220 GHz，第二中频为 40 MHz，因此扫频主路径按 40 MHz 中频进行数字混频。源码中保留的 DDC_IF_HZ 默认值具有历史兼容性质，报告描述以扫频状态机实际设置的 40 MHz 为准。")
    add_doc_paragraph(doc, "DDC 处理时，每个 ADC 采样点先转换为归一化浮点值，再分别乘以 NCO 的 cos 和 -sin 分量，得到同相 I 和正交 Q 两路基带数据。为了保证多次 DMA 累积时相位连续，扫频路径使用独立的 sweep_nco_phase 记录帧间相位推进，而不是每帧都从零相位重新开始。这样可以避免跨帧累积时出现人为相位跳变，提高后续 CIC 抽取和功率估计的一致性。")
    add_doc_paragraph(doc, "与直接对实数中频做 FFT 相比，先做 DDC 的优势在于后续 RBW 滤波可以围绕零频基带进行，滤波器设计和功率估计更加直观。对于扫频频谱仪而言，每个 LO1 频点已经对应一个待测 RF 频率，Zynq 端只需要判断该频点附近 RBW 带宽内的信号能量，因此“DDC 到基带 + 窄带功率估计”的路径比全带 FFT 更贴合本仪器的硬件架构。")

    doc.add_heading("4 RBW滤波与多档分辨率带宽实现", level=2)
    add_doc_paragraph(doc, "分辨率带宽 RBW 是频谱仪最核心的测量参数之一，它决定了相邻频率成分的分辨能力、噪声底显示以及扫频速度。本系统在固件中提供 1 kHz、10 kHz、30 kHz、100 kHz、300 kHz 和 1 MHz 多档 RBW。不同档位不是简单改变显示点数，而是实际改变数字处理链路的抽取比、滤波器长度和观测样本数。这样能够让用户在快速粗扫和窄带精测之间切换，符合频谱仪的典型使用方式。")
    add_doc_paragraph(doc, "RBW 处理链分为两级。第一级为 CIC 抽取器，用较低计算代价完成大倍率降采样；第二级为补偿 FIR 滤波器，在抽取后的较低采样率上形成更接近目标 RBW 的低通响应，并对 CIC 通带下垂和滤波瞬态进行补偿。CIC 的抽取比 R 和级数 N 随 RBW 档位变化，窄 RBW 使用更大的抽取比，从而在较低等效采样率上积累更长时间的样本；宽 RBW 使用较小抽取比，以提高扫频速度。")
    add_doc_paragraph(doc, "补偿 FIR 使用窗函数法动态生成，截止频率由当前 RBW 档位决定，归一化采样率由 ADC_SAMPLE_RATE_HZ / CIC_R 决定。FIR 输出前若直接参与功率计算，会受到初始零状态和滤波器群延迟影响，因此固件为每个 RBW 档位设置 skip points，跳过启动瞬态后再取有效观测窗口。每个频点的累积目标为 observe points、skip points 和 FIR taps 的合计值，并受 ACCUM_BUFFER_SIZE 上限约束。图 2.3.1-4 和表 2.3.1-1 给出了这一处理链和当前参数配置。")
    add_doc_figure(doc, FIGURES[3])
    add_word_table(doc)

    doc.add_heading("5 扫频控制与逐点功率测量", level=2)
    add_doc_paragraph(doc, "扫频计划由 sweep_plan 根据起止频率、中心频率、span 和 RBW 配置生成。当前固件采用与 RBW 绑定的内部步进规则，典型步进约为 RBW / 2，并将最大扫频点数限制为 4096 点。这样做的目的，是避免显示点数固定而导致窄峰被过稀的频率采样跳过；当用户选择更窄的 RBW 时，内部频率步进也随之变小，从而提升窄带信号的捕获能力。")
    add_doc_paragraph(doc, "每个扫频点的状态机流程如图 2.3.1-3 所示。首先根据当前 RF 频点计算并设置 LO1，使输入频率被搬移到第一中频；随后读取锁定指示，等待 PLL 锁定。锁定完成后复位并启动 DMA，采集一帧中频数据；DMA 完成后进入 signal_processing_accumulate_dma，完成 DDC 和 CIC 抽取，并把输出追加到累积缓冲区。如果累积样本不足，则重新启动 DMA；如果已经达到当前 RBW 的累积目标，则进入测量阶段。")
    add_doc_paragraph(doc, "测量阶段先对累积 I/Q 数据应用补偿 FIR，去除瞬态并保留有效观测段，然后对 I/Q 平方和求均值，得到当前频点在 RBW 带宽内的平均功率。完成幅度修正后，sweep_engine 将该点写入内部 points 数组，并调用协议层回调发送频谱点。最后状态机切换到下一频点，复位累积缓冲区和 NCO 连续相位，进入下一轮 LO1 设置。整个流程是非阻塞轮询式实现，便于在裸机主循环中与协议收发、状态查询和性能统计共存。")
    add_doc_figure(doc, FIGURES[2])

    doc.add_heading("6 功率换算与幅度校正", level=2)
    add_doc_paragraph(doc, "当前主测量路径采用时域功率估计，而不是取 FFT 峰值。RBW 滤波后的基带数据包含 I、Q 两路分量，固件对有效观测窗口内所有样本计算 I²+Q² 的平均值，并与 FULL_SCALE_COMPLEX_POWER 进行比较，得到 dBFS。随后加上 ADC_INPUT_FULL_SCALE_DBM，把相对满量程的数字功率换算为 ADC 输入端近似 dBm。该模型与窄带 RBW 滤波相匹配，能够反映当前扫频频点附近一定带宽内的总能量。")
    add_doc_paragraph(doc, "仅有 ADC 输入端功率还不足以代表被测端口的真实输入幅度，因为射频链路中存在 LNA、数字衰减器、VGA、混频链路损耗以及不同频段的增益起伏。固件中的 amplitude_correction 模块根据当前射频前端状态计算总校正量，校正项包括衰减器设置、LNA 增益、VGA 增益、路径校准和频率分段校准。校正后的结果作为 RF 输入端幅度发送给上位机，同时保留 raw power 和 correction dB，便于后续标定和误差分析。")
    add_doc_paragraph(doc, "这一处理方式为低成本硬件提供了可持续改进的幅度定量框架。早期系统可以先使用默认校准项完成相对测量，后续通过标准信号源在多个频点、多个前端状态和多个 RBW 档位下采集误差，再把校准表写入固件或上位机配置。这样不需要改变主处理链路，即可逐步提高整机幅度精度。")
    add_doc_figure(doc, FIGURES[4])

    doc.add_heading("7 数据流式输出与系统优势", level=2)
    add_doc_paragraph(doc, "数字处理完成后，频谱点通过 device_protocol 以 SPECTRUM_DATA 帧流式输出。每个数据点包含当前频率、幅度、总点数、点序号和是否最后一点等信息，上位机接收到后即可逐点更新曲线。相比等待完整扫频结束后一次性上传，流式输出能够缩短用户看到结果的时间，也便于上位机在长 span 或窄 RBW 扫描时显示进度。协议层同时提供 ACK、状态查询和性能 profile 数据，为联调和故障定位提供支撑。")
    add_doc_paragraph(doc, "综合来看，本节数字处理设计有三个突出特点。第一，算法结构与硬件架构一致：硬件负责把宽频输入搬移到固定中频，Zynq 负责稳定的基带化和窄带功率估计。第二，RBW 是真实参与数字滤波和扫频规划的测量参数，而不是单纯的显示选项，因此系统具备频谱仪应有的分辨率带宽概念。第三，扫频状态机、DMA 保护、相位连续 DDC、滤波瞬态跳过和幅度校正共同构成了可调试、可标定、可扩展的工程实现，为后续加入平均检波、峰值保持、零 span、通道功率和自动校准等功能打下基础。")
    add_doc_figure(doc, FIGURES[5], width_cm=15.8)

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run("源码依据：code/pusu_20260516/vitis/pusu_20260516/src 中的 app_config.h、signal_processing.c、sweep_engine.c、sweep_plan.c、dma_capture.c、amplitude_correction.c 和 device_protocol.c。")
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(85, 96, 110)
    doc.save(DOCX_PATH)


def add_slide_title(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(PptInches(0.35), PptInches(0.18), PptInches(12.6), PptInches(0.42))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = PptPt(19)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(38, 50, 56)
    note_box = slide.shapes.add_textbox(PptInches(0.36), PptInches(0.62), PptInches(12.4), PptInches(0.3))
    tf2 = note_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Microsoft YaHei"
    p2.font.size = PptPt(10.5)
    p2.font.color.rgb = PptRGBColor(102, 112, 133)


def generate_ppt():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]
    for meta in FIGURES:
        slide = prs.slides.add_slide(blank)
        add_slide_title(slide, meta["title"], meta["slide_note"])
        image = FIG_DIR / f"{meta['file']}.png"
        with Image.open(image) as im:
            img_w, img_h = im.size
        max_w = 12.05
        max_h = 5.72
        scale = min(max_w / img_w, max_h / img_h)
        pic_w = img_w * scale
        pic_h = img_h * scale
        left = (13.333 - pic_w) / 2
        top = 1.02 + (max_h - pic_h) / 2
        slide.shapes.add_picture(str(image), PptInches(left), PptInches(top), width=PptInches(pic_w), height=PptInches(pic_h))
        cap = slide.shapes.add_textbox(PptInches(0.55), PptInches(6.93), PptInches(12.2), PptInches(0.3))
        tf = cap.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = meta["caption"]
        p.font.name = "Microsoft YaHei"
        p.font.size = PptPt(7.6)
        p.font.color.rgb = PptRGBColor(85, 96, 110)
        p.alignment = PP_ALIGN.CENTER
    prs.save(PPTX_PATH)


def main():
    ensure_dirs()
    generate_figures()
    generate_word()
    generate_ppt()
    print(f"Generated: {DOCX_PATH}")
    print(f"Generated: {PPTX_PATH}")
    print(f"Figures: {FIG_DIR}")


if __name__ == "__main__":
    main()
